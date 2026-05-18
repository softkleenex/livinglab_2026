import collections.abc
from pptx import Presentation
import os

def extract_content(pptx_path):
    prs = Presentation(pptx_path)
    content = []
    
    for i, slide in enumerate(prs.slides):
        slide_data = {
            "slide_number": i + 1,
            "text": [],
            "images": 0
        }
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_data["text"].append(shape.text)
            if shape.shape_type == 13: # Image
                slide_data["images"] += 1
        content.append(slide_data)
    return content

if __name__ == "__main__":
    path = "docs/MDGA_리빙랩_최종발표_2026.pptx"
    try:
        data = extract_content(path)
        for slide in data:
            print(f"--- Slide {slide['slide_number']} (Images: {slide['images']}) ---")
            for t in slide['text']:
                if t.strip():
                    print(t)
            print("\n")
    except Exception as e:
        print(f"Error: {e}")
